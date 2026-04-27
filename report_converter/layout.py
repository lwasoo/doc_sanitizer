from __future__ import annotations

import random
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .common import log, normalize_text, slide_caps
from .models import SlideDraft, TemplateSlide


def extract_template_slides(prs: Presentation) -> list[TemplateSlide]:
    log("读取 PPT 模板结构")
    slides: list[TemplateSlide] = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1 or idx == len(prs.slides):
            continue
        title = ""
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False) and normalize_text(getattr(shp, "text", "")):
                if shp.top <= Inches(1.6):
                    title = normalize_text(shp.text)
                    break
        has_table = any(getattr(shp, "has_table", False) for shp in slide.shapes)
        slides.append(TemplateSlide(idx, title or f"第{idx}页", has_table))
    log(f"PPT 模板解析完成: 需要填充 {len(slides)} 页")
    return slides


def remove_auto_shapes(slide) -> None:
    for shp in list(slide.shapes):
        if shp.name.startswith("AUTO_CONTENT_") or shp.name.startswith("AUTO_LAYOUT_"):
            slide.shapes._spTree.remove(shp._element)


def style_paragraph(paragraph, size: int) -> None:
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(2)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.15


def get_theme_palette(theme: str) -> dict[str, RGBColor]:
    if theme == "corporate_gray":
        return {
            "accent": RGBColor(83, 108, 138),
            "soft": RGBColor(230, 235, 240),
            "line": RGBColor(190, 198, 208),
            "text": RGBColor(48, 72, 96),
        }
    if theme == "legal_red":
        return {
            "accent": RGBColor(161, 47, 47),
            "soft": RGBColor(245, 233, 233),
            "line": RGBColor(219, 188, 188),
            "text": RGBColor(88, 46, 46),
        }
    return {
        "accent": RGBColor(28, 84, 140),
        "soft": RGBColor(232, 240, 248),
        "line": RGBColor(184, 207, 228),
        "text": RGBColor(34, 63, 92),
    }


def suggest_font_size(lines: list[str], base: int = 16) -> int:
    total = sum(len(x) for x in lines)
    longest = max((len(x) for x in lines), default=0)
    size = base
    if total > 260 or longest > 64:
        size -= 1
    if total > 360 or longest > 82:
        size -= 1
    if total > 460 or longest > 100:
        size -= 1
    return max(size, 12)


def _page_budget(title: str, has_table: bool) -> tuple[int, int]:
    bucket = slide_caps(title, has_table)
    max_rows = bucket["page_size"]
    title = normalize_text(title)
    if has_table:
        return 110, max_rows
    if "仲裁" in title or "诉讼" in title or "合同管理" in title or "典型协议" in title:
        return 460, max_rows
    if "概述" in title or "合规" in title or "知识产权" in title:
        return 400, max_rows
    return 420, max_rows


def _row_weight(text: str) -> int:
    text = normalize_text(text)
    if not text:
        return 0
    extra = 0
    if "；" in text or ";" in text:
        extra += 6
    if any(ch.isdigit() for ch in text):
        extra += 4
    return len(text) + extra


def split_balanced(lines: list[str]) -> tuple[list[str], list[str]]:
    if len(lines) <= 1:
        return lines, []
    mid = (len(lines) + 1) // 2
    return lines[:mid], lines[mid:]


def _write_lines_to_box(slide, name: str, left, top, width, height, rows: list[str], font_size: int, bullet: str) -> None:
    if not rows:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = name
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, row in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}{row}" if bullet else row
        style_paragraph(p, size=font_size)
        p.level = 0


def _subtract_interval(intervals: list[tuple[int, int]], block_start: int, block_end: int) -> list[tuple[int, int]]:
    out: list[tuple[int, int]] = []
    for start, end in intervals:
        if block_end <= start or block_start >= end:
            out.append((start, end))
            continue
        if block_start > start:
            out.append((start, block_start))
        if block_end < end:
            out.append((block_end, end))
    return out


def _detect_title_bottom(slide, default_top: int) -> int:
    bottom = default_top
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = normalize_text(getattr(shp, "text", ""))
        if not txt:
            continue
        if shp.top <= Inches(1.6):
            candidate = int(shp.top + shp.height + Inches(0.05))
            if candidate > bottom:
                bottom = candidate
    return bottom


def _find_content_regions(slide, slide_width: int, slide_height: int) -> list[tuple[int, int, int, int]]:
    left = int(Inches(0.55))
    right = int(slide_width - Inches(0.55))
    top_limit = _detect_title_bottom(slide, int(Inches(1.05)))
    bottom_limit = int(slide_height - Inches(0.25))
    intervals: list[tuple[int, int]] = [(top_limit, bottom_limit)]

    for shp in slide.shapes:
        if getattr(shp, "has_table", False):
            t0 = int(max(shp.top - Inches(0.05), top_limit))
            t1 = int(min(shp.top + shp.height + Inches(0.05), bottom_limit))
            intervals = _subtract_interval(intervals, t0, t1)

    intervals = [(start, end) for start, end in intervals if end - start >= int(Inches(0.45))]
    if not intervals:
        return [(left, top_limit, right - left, int(max(bottom_limit - top_limit, Inches(0.6))))]
    return [(left, start, right - left, end - start) for start, end in intervals]


def _find_content_region(slide, slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    regions = _find_content_regions(slide, slide_width, slide_height)
    return max(regions, key=lambda item: item[3])


def _write_rows_across_regions(
    slide,
    base_name: str,
    rows: list[str],
    regions: list[tuple[int, int, int, int]],
    font_size: int,
    bullet: str,
) -> None:
    if not rows or not regions:
        return

    total_height = sum(region[3] for region in regions)
    start = 0
    remaining_rows = len(rows)
    remaining_regions = len(regions)

    for idx, (left, top, width, height) in enumerate(regions):
        if remaining_rows <= 0:
            break
        if idx == len(regions) - 1:
            chunk = rows[start:]
        else:
            ratio = height / max(total_height, 1)
            take = max(1, round(len(rows) * ratio))
            take = min(take, remaining_rows - (remaining_regions - 1))
            chunk = rows[start : start + take]
        _write_lines_to_box(slide, f"{base_name}_{idx}", left, top, width, height, chunk, font_size, bullet)
        start += len(chunk)
        remaining_rows -= len(chunk)
        remaining_regions -= 1


def add_content_textbox(slide, draft: SlideDraft, has_table: bool, slide_width: int, slide_height: int) -> None:
    if not draft.bullets:
        return

    rows = draft.bullets
    font_size = suggest_font_size(rows, base=15 if has_table else 16)
    if has_table:
        regions = _find_content_regions(slide, slide_width, slide_height)
        _write_rows_across_regions(slide, f"AUTO_CONTENT_{draft.slide_index}", rows, regions, font_size, "• ")
        return

    left, top, width, height = _find_content_region(slide, slide_width, slide_height)
    if height < int(Inches(0.45)):
        return
    _write_lines_to_box(slide, f"AUTO_CONTENT_{draft.slide_index}", left, top, width, height, rows, font_size, bullet="• ")


def choose_formal_variant(title: str, bullets: list[str], diversity: str, seed: int, slide_index: int) -> str:
    title = title or ""
    if diversity == "none":
        return "single_column"
    if len(bullets) >= 6:
        return "two_column"
    if any(k in title for k in ["风险", "诉讼", "仲裁"]):
        base = ["risk_matrix", "timeline", "two_column"]
    elif any(k in title for k in ["数据", "统计", "流程"]):
        base = ["kpi_cards", "two_column", "single_column"]
    elif any(k in title for k in ["项目", "进度", "政策", "合规"]):
        base = ["timeline", "two_column", "single_column"]
    else:
        base = ["two_column", "single_column", "kpi_cards"]

    if len(bullets) <= 2:
        return "single_column"
    if len(bullets) >= 4 and diversity == "high":
        base = ["two_column", "timeline", "kpi_cards", "risk_matrix"]
    elif diversity == "low":
        base = [base[0]]

    rng = random.Random((seed or 0) * 131 + slide_index * 17 + len(bullets))
    return base[rng.randrange(len(base))]


def add_visual_accent(slide, slide_index: int, palette: dict[str, RGBColor]) -> None:
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(1.16), Inches(0.06), Inches(3.9))
    accent.name = f"AUTO_LAYOUT_ACCENT_{slide_index}"
    accent.fill.solid()
    accent.fill.fore_color.rgb = palette["accent"]
    accent.line.fill.background()


def render_single_column(slide, slide_index: int, rows: list[str], font_size: int) -> None:
    _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}", Inches(0.72), Inches(1.12), Inches(8.55), Inches(5.75), rows, font_size, bullet="• ")


def render_two_column(slide, slide_index: int, rows: list[str], font_size: int) -> None:
    left_rows, right_rows = split_balanced(rows)
    col_top = Inches(1.12)
    col_height = Inches(5.75)
    col_width = Inches(4.18)
    _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_L", Inches(0.72), col_top, col_width, col_height, left_rows, font_size, bullet="• ")
    _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_R", Inches(5.08), col_top, col_width, col_height, right_rows, font_size, bullet="• ")


def render_timeline(slide, slide_index: int, rows: list[str], font_size: int, palette: dict[str, RGBColor]) -> None:
    start_top = Inches(1.18)
    step = Inches(0.96)
    for idx, row in enumerate(rows[:6]):
        y = start_top + idx * step
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.92), y + Inches(0.05), Inches(0.12), Inches(0.12))
        dot.name = f"AUTO_LAYOUT_DOT_{slide_index}_{idx}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = palette["accent"]
        dot.line.fill.background()
        if idx < len(rows) - 1:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.97), y + Inches(0.16), Inches(0.02), Inches(0.6))
            line.name = f"AUTO_LAYOUT_LINE_{slide_index}_{idx}"
            line.fill.solid()
            line.fill.fore_color.rgb = palette["line"]
            line.line.fill.background()
        _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_{idx}", Inches(1.2), y, Inches(7.9), Inches(0.68), [row], font_size, bullet="")


def render_kpi_cards(slide, slide_index: int, rows: list[str], font_size: int, palette: dict[str, RGBColor]) -> None:
    rows = rows[:4]
    card_w = Inches(4.1)
    card_h = Inches(1.0)
    positions = [(Inches(0.72), Inches(1.22)), (Inches(4.95), Inches(1.22)), (Inches(0.72), Inches(2.45)), (Inches(4.95), Inches(2.45))]
    for idx, row in enumerate(rows):
        left, top = positions[idx]
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.name = f"AUTO_LAYOUT_CARD_{slide_index}_{idx}"
        card.fill.solid()
        card.fill.fore_color.rgb = palette["soft"]
        card.line.color.rgb = palette["line"]
        _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_{idx}", left + Inches(0.18), top + Inches(0.12), card_w - Inches(0.32), card_h - Inches(0.24), [row], font_size, bullet="")


def render_risk_matrix(slide, slide_index: int, rows: list[str], font_size: int, palette: dict[str, RGBColor]) -> None:
    left = Inches(0.86)
    top = Inches(1.24)
    width = Inches(8.1)
    height = Inches(5.2)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    box.name = f"AUTO_LAYOUT_RISKBOX_{slide_index}"
    box.fill.solid()
    box.fill.fore_color.rgb = palette["soft"]
    box.line.color.rgb = palette["line"]
    for idx, row in enumerate(rows[:6]):
        _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_{idx}", left + Inches(0.22), top + Inches(0.18 + idx * 0.82), width - Inches(0.44), Inches(0.62), [row], font_size, bullet="• ")


def add_formal_layout_content_v2(
    slide,
    draft: SlideDraft,
    has_table: bool,
    title: str,
    theme: str,
    diversity: str,
    seed: int,
    slide_width: int | None = None,
    slide_height: int | None = None,
) -> bool:
    if not draft.bullets:
        return False

    if has_table:
        if slide_width is None or slide_height is None:
            slide_width = int(Inches(10))
            slide_height = int(Inches(7.5))
        rows = draft.bullets
        font_size = suggest_font_size(rows, base=13)
        regions = _find_content_regions(slide, slide_width, slide_height)
        _write_rows_across_regions(slide, f"AUTO_CONTENT_{draft.slide_index}", rows, regions, font_size, "• ")
        return True

    palette = get_theme_palette(theme)
    rows = draft.bullets

    add_visual_accent(slide, draft.slide_index, palette)
    font_size = suggest_font_size(rows, base=15)
    variant = choose_formal_variant(title, rows, diversity, seed, draft.slide_index)

    if variant == "two_column":
        render_two_column(slide, draft.slide_index, rows, font_size)
    elif variant == "timeline":
        render_timeline(slide, draft.slide_index, rows, font_size, palette)
    elif variant == "kpi_cards":
        render_kpi_cards(slide, draft.slide_index, rows, font_size, palette)
    elif variant == "risk_matrix":
        render_risk_matrix(slide, draft.slide_index, rows, font_size, palette)
    else:
        render_single_column(slide, draft.slide_index, rows, font_size)
    return False


def fill_table_metrics(slide, metrics: dict[str, str]) -> None:
    alias_map = {
        "专利申请量": ["专利申请量", "其他知识产权申请", "BU10申请量", "BU11申请量", "BU16申请量"],
        "专利调查量": ["专利调查量", "专利调查量BU10", "专利调查量BU11", "专利调查量BU16"],
        "一般文件用印": ["一般文件用印"],
        "法律文件用印": ["法律文件用印"],
        "集团制式文件用印": ["集团制式文件用印"],
        "非制式文件-供应商": ["非制式文件-供应商"],
        "非制式文件-客户": ["非制式文件-客户"],
        "非制式文件-内部行政": ["非制式文件-内部行政"],
        "非制式文件-重要文件": ["非制式文件-重要文件"],
    }

    def resolve_metric(label: str) -> str:
        norm = normalize_text(label)
        if norm in metrics:
            return metrics.get(norm, "-")
        for key, cands in alias_map.items():
            if norm == key or norm in cands:
                for cand in cands:
                    val = metrics.get(cand)
                    if val and val != "-":
                        return val
        if "申请" in norm and "专利" in norm:
            return metrics.get("专利申请量") or metrics.get("其他知识产权申请") or "-"
        if "调查" in norm and "专利" in norm:
            return metrics.get("专利调查量") or metrics.get("专利调查量BU10") or "-"
        if "一般文件" in norm:
            return metrics.get("一般文件用印", "-")
        if "法律文件" in norm:
            return metrics.get("法律文件用印", "-")
        if "集团制式" in norm:
            return metrics.get("集团制式文件用印", "-")
        if "供应商" in norm and "非制式" in norm:
            return metrics.get("非制式文件-供应商", "-")
        return "-"

    for shp in slide.shapes:
        if not getattr(shp, "has_table", False):
            continue
        table = shp.table
        if len(table.columns) >= 3:
            header_apply = normalize_text(table.cell(0, 1).text)
            header_survey = normalize_text(table.cell(0, 2).text)
            for row_idx in range(1, len(table.rows)):
                row_label = normalize_text(table.cell(row_idx, 0).text).replace(" ", "").upper()
                if row_label in {"BU10", "BU11", "BU16"}:
                    table.cell(row_idx, 1).text = metrics.get(f"{row_label}申请量", metrics.get("专利申请量", "-")) if "申请" in header_apply else "-"
                    table.cell(row_idx, 2).text = metrics.get(f"专利调查量{row_label}", metrics.get("专利调查量", "-")) if "调查" in header_survey else "-"
            continue
        for row_idx in range(1, len(table.rows)):
            label = normalize_text(table.cell(row_idx, 0).text)
            table.cell(row_idx, 1).text = resolve_metric(label)


def split_into_pages(bullets: list[str], title: str, has_table: bool) -> list[list[str]]:
    if not bullets:
        return [[]]
    budget, max_rows = _page_budget(title, has_table)
    target_fill = int(budget * 0.68)
    pages: list[list[str]] = []
    current: list[str] = []
    used = 0

    for bullet in bullets:
        weight = _row_weight(bullet)
        if not current:
            current = [bullet]
            used = weight
            continue

        over_rows = len(current) >= max_rows
        over_budget = used + weight > budget
        enough_fill = used >= target_fill

        if over_rows or (over_budget and enough_fill):
            pages.append(current)
            current = [bullet]
            used = weight
            continue

        current.append(bullet)
        used += weight

    if current:
        pages.append(current)
    return pages


def insert_slide_after(prs: Presentation, after_index: int, layout) -> Any:
    new_slide = prs.slides.add_slide(layout)
    sld_id_lst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(after_index + 1, new_id)
    return prs.slides[after_index + 1]


def set_slide_title_text(slide, title: str) -> None:
    target = None
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = normalize_text(getattr(shp, "text", ""))
        if shp.top <= Inches(1.6):
            target = shp
            if txt:
                break
    if target is None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(8.5), Inches(0.6))
        box.name = "AUTO_OVERFLOW_TITLE"
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        return
    target.text = title
